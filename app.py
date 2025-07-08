# app.py – SK‑StratAI slide generator
# Streamlit • LangChain 0.2 • OpenAI 4o-mini • SerpAPI • PPTX (Calibri 12 pt)
# ============================================================================
"""
Chem‑Energy‑Materials market‑context slide builder for Simon‑Kucher consultants.

Key features
• Fixed GPT model (gpt-4o-mini)
• Typo‑tolerant client-name suggestions via SerpAPI
• Domain selection from live web results
• Ranked multiselect verticals based on website & product lines
• Regulation snippets (last 2 yrs)
• 10‑20 MECE bullets per slide
• Inline citations as hidden hyperlinks
• Calibri 12 pt bullets
• Non‑blocking follow‑up Qs
• Robust JSON parsing: double-quoted keys/values, no single quotes
"""

# -------------------------------------------------------------------------
# Imports & constants
# -------------------------------------------------------------------------
import os, json, re, datetime, tempfile, pathlib, requests, difflib
from typing import List

import streamlit as st
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from serpapi import GoogleSearch
import openai

# Force use of gpt-4o-mini
MODEL_NAME = "gpt-4o-mini"

EXT_TAXONOMY = [
    "Commodity Chemicals", "Specialty Chemicals", "Petrochemicals & Refining",
    "Industrial Gases", "Energy Storage Materials", "Battery Materials",
    "Advanced Materials & Additives", "Composites", "3D-print Powders",
    "Crop Protection", "Bio-control", "Seeds & Traits", "Pest Control",
    "Bio-based Systems", "Fermentation Ingredients", "Enzymes & Biologics",
    "Circular Economy", "Plastic Recycling", "Water Treatment Chemicals",
    "Paints & Coatings", "Construction Chemicals", "Personal Care Ingredients",
    "Agro Bio-stimulants", "Fertilizers", "Cement", "Glass & Ceramics",
    "Steel & Metallurgy", "Rubber & Elastomers", "Catalysts"
]
EMBEDDER = OpenAIEmbeddings(model="text-embedding-3-small")
THIS_YEAR = datetime.datetime.now().year

# -------------------------------------------------------------------------
# SerpAPI helpers
# -------------------------------------------------------------------------
@st.cache_data(show_spinner=False)
def serp(query: str, k: int = 5) -> dict:
    key = os.getenv("SERPAPI_API_KEY")
    if not key:
        return {}
    try:
        return GoogleSearch({"q": query, "api_key": key, "num": k, "engine": "google"}).get_dict()
    except Exception:
        return {}

@st.cache_data(show_spinner=False)
def serp_snippets(query: str, k: int = 5) -> List[str]:
    return [r.get("snippet", "")[:500] for r in serp(query, k).get("organic_results", [])[:k]]

@st.cache_data(show_spinner=False)
def spelling_fix(name: str) -> str | None:
    info = serp(name, 3).get("search_information", {})
    return info.get("spelling_fix")

@st.cache_data(show_spinner=False)
def suggest_client_names(name: str, k: int = 5) -> List[str]:
    data = serp(name, k)
    titles = []
    for r in data.get("organic_results", []):
        title = r.get("title") or r.get("link")
        if title:
            titles.append(title)
    return titles

@st.cache_data(show_spinner=False)
def find_corporate_domain(name: str) -> str | None:
    data = serp(name, 5)
    for res in data.get("organic_results", []):
        url = res.get("link", "")
        m = re.match(r"https?://(?:www\.)?([\w.\-]+\.[a-z]{2,})", url)
        if m:
            dom = m.group(1).lower()
            if not any(bad in dom for bad in ("news.", "linkedin", "wikipedia")):
                return dom
    guess = name.lower().replace(" ", "") + ".com"
    try:
        requests.head(f"https://{guess}", timeout=4).raise_for_status()
        return guess
    except Exception:
        return None

# -------------------------------------------------------------------------
# Scrape homepage text
# -------------------------------------------------------------------------
@st.cache_data(show_spinner=False)
def scrape_home(domain: str, max_chars: int = 8000) -> str:
    if not domain:
        return ""
    try:
        html = requests.get(f"https://{domain}", timeout=8).text
        soup = BeautifulSoup(html, "html.parser")
        return re.sub(r"\s+", " ", soup.get_text(" ", strip=True))[:max_chars]
    except Exception:
        return ""

# -------------------------------------------------------------------------
# Vertical suggestion via embeddings
# -------------------------------------------------------------------------
@st.cache_data(show_spinner=False)
def suggest_verticals(text: str, top_n: int = 5) -> List[str]:
    if not text.strip():
        return []
    try:
        vec = EMBEDDER.embed_query(text[:2048])
    except Exception:
        return []
    scores = [(sum(a*b for a,b in zip(vec, EMBEDDER.embed_query(label))), label)
              for label in EXT_TAXONOMY]
    return [lbl for _, lbl in sorted(scores, reverse=True)[:top_n]]

# -------------------------------------------------------------------------
# Streamlit UI
# -------------------------------------------------------------------------
st.set_page_config(page_title="SK‑StratAI", layout="centered")
st.header("Generate market-context slides")

input_name = st.text_input("Client name (required)")
# Spelling & name suggestions
candidates = []
if input_name:
    fix = spelling_fix(input_name)
    if fix and fix.lower() != input_name.lower():
        candidates.append(fix)
    candidates += suggest_client_names(input_name)
candidates = list(dict.fromkeys(candidates))
selection = st.selectbox("Did you mean…?", [input_name] + candidates) if candidates else input_name
client = selection

# Domain selection & debug
candidates_raw = serp(client, 5)
with st.expander("🔍 SerpAPI domain search debug"):
    if not os.getenv("SERPAPI_API_KEY"):
        st.error("SERPAPI_API_KEY not set or invalid. Please configure your key.")
    st.json(candidates_raw)
domains = [re.match(r"https?://(?:www\.)?([\w.\-]+\.[a-z]{2,})", r.get("link",""))
           for r in candidates_raw.get("organic_results", [])]
domains = [m.group(1) for m in domains if m]
selected_domain = st.selectbox("Select domain", domains or [client.lower().replace(" ","")+".com"] )
home_text = scrape_home(selected_domain)

# Industry picker
quick = serp_snippets(f"{client} product lines", k=3)
verticals = suggest_verticals(home_text + " ".join(quick)) or EXT_TAXONOMY[:5]
opts = verticals + [x for x in EXT_TAXONOMY if x not in verticals]
selected = st.multiselect("Relevant verticals", opts, default=verticals)
industry = ", ".join(selected)

brief = st.text_area("Short client brief", height=100)
uploads = st.file_uploader("Upload decks (optional)", type=["pptx"], accept_multiple_files=True)
if not st.button("Generate slides"): st.stop()

# -------------------------------------------------------------------------
# Build context & prompt
# -------------------------------------------------------------------------
retrieved: List[dict] = []
if pathlib.Path("vector_db").exists():
    db = Chroma(persist_directory="vector_db", embedding_function=EMBEDDER)
    retrieved += db.similarity_search(brief or industry or client, k=8)
if uploads:
    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    for up in uploads:
        tf = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        tf.write(up.read()); tf.close()
        docs = splitter.split_documents(UnstructuredPowerPointLoader(tf.name).load())
        retrieved += docs
webs = serp_snippets(f"{client} {industry} market trend {THIS_YEAR}", k=5)
retrieved += [{"page_content":s, "metadata":{"source":"web"}} for s in webs]
regs = []
for yr in (THIS_YEAR, THIS_YEAR-1):
    regs += serp_snippets(f"{industry} regulation {yr} pricing OR go-to-market OR channel", k=3)
retrieved += [{"page_content":s, "metadata":{"source":"reg"}} for s in regs]
blocks = [f"SRC:{d.get('metadata',{}).get('source','internal')}|{d.get('page_content','')[:500]}"
          for d in retrieved]
blocks = blocks[:20]

SYSTEM_PROMPT = '''
You are **SK‑StratAI**, a senior virtual consultant at **Simon‑Kucher & Partners (SKP)** with deep expertise in Growth Strategy, Go‑to‑Market (GTM), Pricing Strategy & Price Setting, Commercial & Marketing Excellence, Sales Acceleration, and Product‑Market‑Fit. Adopt the crisp, insight‑dense style of an SKP Partner.

1. Output format  
Return a single, valid JSON object containing exactly three mandatory arrays and one optional array:

{
  "market_slide":   [ /* 10‑20 strings */ ],
  "client_slide":   [ /* 10‑20 strings */ ],
  "question_slide": [ /* 20    strings */ ],
  "follow_up":      [ /* ≤5 strings; include ONLY when critical data are missing */ ]
}

Do not wrap the JSON in markdown fences and do not add any other keys or text.

2. Array instructions

Array: market_slide  
• 10‑20 MECE, consultant‑style bullets on the most material external trends, pitfalls, and challenges in the listed competence areas.  
• Must reference noteworthy competitor moves and cite ≥1 explicitly named regulation enacted in the last 24 months that affects GTM, pricing, channel, or growth.

Array: client_slide  
• 10‑20 MECE, insight‑driven bullets that map the issues from market_slide to the client’s current positioning, capabilities, and gaps.

Array: question_slide  
• Exactly 20 sharp, Partner‑level diagnostic questions that surface latent needs and quantifiable value levers across Growth, GTM, Pricing, Commercial Excellence, Sales Acceleration, and Product‑Market‑Fit.  
• Total combined length ≤100 words.

Optional array: follow_up  
• Up to 5 concise questions requesting only the missing information strictly needed to meet the specs above.

3. Writing & citation rules  
• Tone: Authoritative, concise SKP consulting voice. No emojis.  
• Freshness: Base insights on sources published ≤24 months ago.  
• Citations: Where a bullet relies on an external fact, append a superscript id in parentheses, e.g. (¹), and embed the source URL in that id as a hyperlink.  
• MECE: Ensure bullets are mutually exclusive and collectively exhaustive.  
• Data gaps: If essential data are unavailable, still populate the three core arrays and add follow_up.

Deliver only the JSON object; no prose commentary. Use only double quotes for keys and string values.


'''

payload = {"client":client, "industry":industry, "brief":brief, "context":"\n\n".join(blocks) or "none"}
msgs=[{"role":"system","content":SYSTEM_PROMPT},{"role":"user","content":json.dumps(payload)}]

# -------------------------------------------------------------------------
# OpenAI call & parse
# -------------------------------------------------------------------------
def call(**kw):
    try: return openai.chat.completions.create(**kw)
    except Exception as e: return e
base=dict(model=MODEL_NAME, messages=msgs, temperature=0.3)
res = call(**base, response_format={"type":"json_object"}, max_tokens=900)
if isinstance(res, Exception) and "response_format" in str(res):
    res = call(**base, max_tokens=900)
if isinstance(res, Exception) and "max_tokens" in str(res):
    res = call(**base, max_completion_tokens=900)
if isinstance(res, Exception):
    st.error(f"API err {res}")
    st.stop()

raw=res.choices[0].message.content.strip()
if raw.startswith("```json"): raw=re.sub(r"^```json\s*|\s*```$","",raw,flags=re.S).strip()
m=re.search(r"\{.*\}",raw,flags=re.S); raw=m.group(0) if m else raw
raw=re.sub(r",\s*([}\]])",r"\1",raw)
raw=re.sub(r'"\s*\n\s*([^"].*?)\s*\n\s*"',lambda m:'"'+m.group(1).replace("\n"," ")+'"',raw)
try: data=json.loads(raw)
except Exception as e: st.error(f"JSON err → {e}"); st.text(raw[:500]); st.stop()
if data.get("follow_up"):
    with st.expander("Follow-up"):
        for q in data["follow_up"]:
            st.write(f"• {q}")

# -------------------------------------------------------------------------
# Build PPT
# -------------------------------------------------------------------------
prs=Presentation()
def slide(t,b):
    s=prs.slides.add_slide(prs.slide_layouts[1]); s.shapes.title.text=t
    tf=s.shapes.placeholders[1].text_frame; tf.clear()
    for btxt in b:
        m=re.search(r"\((\d+)\)\s*<(?P<url>https?://\S+)>",btxt)
        vis=re.sub(r"\s*<https?://\S+>","",btxt)
        p=tf.add_paragraph(); p.text=vis; p.font.name="Calibri"; p.font.size=Pt(12)
        if m:
            cid,url=m.group(1),m.group("url")
            for run in p.runs:
                if f"({cid})" in run.text:
                    run.font.color.rgb=RGBColor(0,102,204); run._r.hyperlink.address=url
            s.notes_slide.notes_text_frame.add_paragraph().text=f"({cid}) {url}"
slide("Market context",data.get("market_slide",[]))
slide("Client context",data.get("client_slide",[]))
slide("Questions",data.get("question_slide",[]))
fn=f"slides-{client.replace(' ','_')}-{datetime.datetime.now():%Y%m%d-%H%M}.pptx"
prs.save(fn)
with open(fn,'rb') as f: st.download_button("Download PPT",f.read(),file_name=fn,mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
